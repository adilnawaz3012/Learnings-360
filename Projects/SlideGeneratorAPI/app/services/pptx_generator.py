from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from ..models.schemas import PresentationConfig, SlideType, AspectRatio
from ..utils.template_loader import get_template
from concurrent.futures import ThreadPoolExecutor

def create_pptx(config: PresentationConfig, template_id: str = "modern", 
                aspect_ratio: AspectRatio = AspectRatio.WIDESCREEN, 
                custom_size: tuple = None) -> BytesIO:
    """Create PowerPoint presentation with template and aspect ratio"""
    template = get_template(template_id)
    prs = Presentation()
    
    # Apply aspect ratio
    apply_aspect_ratio(prs, aspect_ratio, custom_size)
    
    # Apply template styles
    apply_template_styles(prs, template)
    
    # Add slides
    for slide_config in config.slides:
        add_slide(prs, slide_config)
    
    # Save to bytes buffer
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

def apply_aspect_ratio(prs, aspect_ratio, custom_size):
    if aspect_ratio == AspectRatio.STANDARD:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    elif aspect_ratio == AspectRatio.WIDESCREEN:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    elif aspect_ratio == AspectRatio.CUSTOM and custom_size:
        prs.slide_width = Inches(custom_size[0])
        prs.slide_height = Inches(custom_size[1])

def apply_template_styles(prs, template):
    # Apply template-specific styling
    # This would be implemented based on template configuration
    pass

def add_slide(prs, slide_config):
    # Add slide based on type
    if slide_config.slide_type == SlideType.TITLE:
        add_title_slide(prs, slide_config)
    elif slide_config.slide_type == SlideType.TITLE_AND_CONTENT:
        add_title_content_slide(prs, slide_config)
    # ... other slide types

def add_title_slide(prs, slide_config):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = slide_config.title

def add_title_content_slide(prs, slide_config):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = slide_config.title
    content.text = slide_config.content or ""
    
    # Add image if available
    if slide_config.image_url:
        # Implementation to add image
        pass