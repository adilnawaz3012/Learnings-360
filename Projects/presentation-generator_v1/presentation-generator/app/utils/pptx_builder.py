# app/utils/pptx_builder.py
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from app.models.presentation_models import PresentationData, Slide, SlideLayout
from app.core.config import logger

def create_presentation_file(data: PresentationData, config) -> str:
    """Generates a .pptx file from presentation data and returns the file path."""
    prs = PptxPresentation()
    
    # Apply a very basic theme based on config
    is_dark_theme = config.theme == 'dark'
    bg_color = RGBColor(0x20, 0x21, 0x23) if is_dark_theme else RGBColor(0xFF, 0xFF, 0xFF)
    text_color = RGBColor(0xFF, 0xFF, 0xFF) if is_dark_theme else RGBColor(0x00, 0x00, 0x00)

    for slide_data in data.slides:
        slide_layout_name = _get_pptx_layout(slide_data.type)
        slide_layout = prs.slide_layouts[slide_layout_name]
        slide = prs.slides.add_slide(slide_layout)

        # Set slide background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Populate slide content
        _populate_slide_content(slide, slide_data, text_color)

    file_path = f"generated_presentations/{data.title.replace(' ', '_').lower()}.pptx"
    prs.save(file_path)
    logger.info(f"Presentation saved to {file_path}")
    return file_path

def _get_pptx_layout(layout_type: SlideLayout) -> int:
    """Maps our custom layout type to python-pptx's default layout indices."""
    if layout_type == SlideLayout.TITLE:
        return 0  # Title Slide
    if layout_type == SlideLayout.TWO_COLUMN:
        return 3 # Two Content
    # Use 'Title and Content' for bullets and content+image
    return 1

def _populate_slide_content(slide, data: Slide, text_color: RGBColor):
    """Fills a slide with content based on its type."""
    if data.title:
        title_shape = slide.shapes.title
        title_shape.text = data.title
        title_shape.text_frame.paragraphs[0].font.color.rgb = text_color

    # Using placeholders if they exist
    body_shape = None
    if slide.placeholders:
        for shape in slide.placeholders:
            if shape.name.startswith(('Content Placeholder', 'Text Placeholder', 'Body')):
                 body_shape = shape
                 break
    
    if data.type == SlideLayout.TITLE and data.subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = data.subtitle
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = text_color
    
    elif data.type == SlideLayout.BULLET_POINTS and body_shape:
        tf = body_shape.text_frame
        tf.clear()
        for point in data.points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.color.rgb = text_color

    elif data.type == SlideLayout.CONTENT_WITH_IMAGE and body_shape:
        body_shape.text = data.content
        body_shape.text_frame.paragraphs[0].font.color.rgb = text_color
        # Add a placeholder note for the image
        slide.shapes.add_textbox(Inches(6), Inches(2), Inches(3.5), Inches(3.5)).text = f"[Image: {data.image_suggestion}]"
        
    elif data.type == SlideLayout.TWO_COLUMN:
        # Assumes layout 3 'Two Content'
        left_ph = slide.placeholders[1]
        right_ph = slide.placeholders[2]
        left_ph.text = data.left_content
        right_ph.text = data.right_content
        left_ph.text_frame.paragraphs[0].font.color.rgb = text_color
        right_ph.text_frame.paragraphs[0].font.color.rgb = text_color