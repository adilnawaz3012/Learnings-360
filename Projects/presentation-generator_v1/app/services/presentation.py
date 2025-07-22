from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import logging
from app.models.presentation import SlideConfig, ThemeConfig, SlideLayout, AspectRatio
from app.core.config import settings
from app.utils.redis import cache
from app.core.exceptions import PPTXGenerationException
from concurrent.futures import ThreadPoolExecutor, as_completed

logger = logging.getLogger(__name__)

class PPTXGenerator:
    def __init__(self, theme: ThemeConfig):
        self.theme = theme
        self.aspect_ratio = self._get_aspect_ratio(theme.aspect_ratio)
    
    def _get_aspect_ratio(self, ratio: str):
        return settings.ASPECT_RATIOS.get(ratio, settings.ASPECT_RATIOS[settings.DEFAULT_ASPECT_RATIO])
    
    def create_presentation(self, slides: list[SlideConfig], output_path: str):
        try:
            template_path = os.path.join(settings.TEMPLATE_DIR, self.theme.template or settings.DEFAULT_TEMPLATE)
            
            if os.path.exists(template_path):
                prs = Presentation(template_path)
            else:
                prs = Presentation()
                width, height = self.aspect_ratio
                prs.slide_width = Inches(width)
                prs.slide_height = Inches(height)
            
            self._apply_theme_to_master(prs)
            
            with ThreadPoolExecutor(max_workers=min(8, len(slides))) as executor:
                futures = {
                    executor.submit(self._add_slide, prs, slide): idx
                    for idx, slide in enumerate(slides)
                }
                
                for future in as_completed(futures):
                    try:
                        future.result()
                    except Exception as e:
                        idx = futures[future]
                        logger.error(f"Failed to generate slide {idx}: {str(e)}")
            
            prs.save(output_path)
            return output_path
        except Exception as e:
            logger.error(f"PPTX generation failed: {str(e)}")
            raise PPTXGenerationException(f"Failed to create presentation: {str(e)}")

    def _apply_theme_to_master(self, prs):
        slide_master = prs.slide_master
        background = slide_master.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor.from_string(self.theme.secondary_color)
        
        title_style = slide_master.placeholders[0].text_frame.paragraphs[0].font
        title_style.name = self.theme.font
        title_style.size = Pt(self.theme.title_font_size)
        title_style.color.rgb = RGBColor.from_string(self.theme.primary_color)
        
        for placeholder in slide_master.placeholders[1:]:
            for paragraph in placeholder.text_frame.paragraphs:
                paragraph.font.name = self.theme.font
                paragraph.font.size = Pt(self.theme.content_font_size)
                paragraph.font.color.rgb = RGBColor.from_string(self.theme.primary_color)

    def _add_slide(self, prs, slide_cfg: SlideConfig):
        layout_map = {
            SlideLayout.TITLE: 0,
            SlideLayout.BULLET_POINTS: 1,
            SlideLayout.TWO_COLUMN: 3,
            SlideLayout.IMAGE_PLACEHOLDER: 7
        }
        
        slide_layout = prs.slide_layouts[layout_map[slide_cfg.layout]]
        slide = prs.slides.add_slide(slide_layout)
        content = slide_cfg.content
        
        title = slide.shapes.title
        title.text = content.title
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.theme.accent_color)
        
        placeholders = slide.shapes.placeholders
        
        if slide_cfg.layout == SlideLayout.TITLE and len(placeholders) > 1:
            subtitle = placeholders[1]
            subtitle.text = content.bullets[0] if content.bullets else ""

        elif slide_cfg.layout == SlideLayout.BULLET_POINTS and len(placeholders) > 1:
            content_placeholder = placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for bullet in content.bullets or []:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.name = self.theme.font
                p.font.size = Pt(self.theme.content_font_size)
                p.font.color.rgb = RGBColor.from_string(self.theme.primary_color)

        elif slide_cfg.layout == SlideLayout.TWO_COLUMN and len(placeholders) > 2:
            left_col = placeholders[1]
            right_col = placeholders[2]
            
            if content.columns:
                left_col.text = content.columns.get("left", "")
                right_col.text = content.columns.get("right", "")
            elif content.bullets:
                midpoint = len(content.bullets) // 2
                left_col.text = "\n".join(content.bullets[:midpoint])
                right_col.text = "\n".join(content.bullets[midpoint:])

        elif slide_cfg.layout == SlideLayout.IMAGE_PLACEHOLDER and len(placeholders) > 2:
            text_placeholder = placeholders[1]
            image_placeholder = placeholders[2]
            
            if content.bullets:
                text_placeholder.text = "\n".join(content.bullets)
            
            if content.image_suggestion:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = f"Image suggestion: {content.image_suggestion}"
        
        if content.citations:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(prs.slide_height.inches - 1), 
                                           Inches(prs.slide_width.inches - 1), Inches(0.5))
            tf = txBox.text_frame
            tf.text = "Sources: " + "; ".join(content.citations)
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = RGBColor.from_string("#666666")

class PresentationService:
    def __init__(self, repo):
        self.repo = repo
    
    async def generate_pptx(self, presentation):
        cache_key = f"pptx:{presentation.id}"
        if cached_path := cache.get(cache_key):
            if os.path.exists(cached_path):
                return cached_path
        
        lock_key = f"lock:pptx:{presentation.id}"
        if not cache.acquire_lock(lock_key):
            raise PPTXGenerationException("Presentation is being generated by another request")
        
        try:
            if cached_path := cache.get(cache_key):
                if os.path.exists(cached_path):
                    return cached_path
            
            output_path = os.path.join(settings.STORAGE_PATH, f"{presentation.id}.pptx")
            generator = PPTXGenerator(presentation.theme)
            generator.create_presentation(presentation.slides, output_path)
            
            cache.set(cache_key, output_path)
            return output_path
        finally:
            cache.release_lock(lock_key)