from pptx import Presentation
from pptx.util import Inches
import os

ASPECT_RATIO_DIMENSIONS = {
    "16:9": (13.33, 7.5),
    "4:3": (10.0, 7.5)
}

def create_presentation(presentation_id, content, request):
    template_path = "templates/default_template.pptx"
    prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()

    width, height = ASPECT_RATIO_DIMENSIONS.get(request.aspect_ratio or "16:9", (13.33, 7.5))
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    for slide_data in content["slides"]:
        if slide_data["layout"] == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = slide_data["title"]
            slide.placeholders[1].text = slide_data.get("content", "")
        elif slide_data["layout"] == "bullet":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data["title"]
            body = slide.placeholders[1].text_frame
            for point in slide_data["points"]:
                body.add_paragraph().text = point

    path = f"/app/generated/{presentation_id}.pptx"
    os.makedirs("/app/generated", exist_ok=True)
    prs.save(path)
    return path

def get_presentation_path(path):
    from fastapi.responses import FileResponse
    return FileResponse(path, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename=os.path.basename(path))